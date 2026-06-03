"""In-place anonymisering av originalfiler (DOCX, PPTX, XLSX, PDF).

Strategi per format:
  DOCX  – python-docx: merge runs per avsnitt, erstatt tekst, skriv tilbake
  PPTX  – python-pptx: samme run-merge per tekstramme
  XLSX  – openpyxl:    erstatt celleverdier direkte
  PDF   – pymupdf:     søk+rediger med standard PDF-redaksjonsannotering
"""
from __future__ import annotations

import unicodedata
import re
import os
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

from xlent_scanner.anonymize import _clean_replacement_text, _replace_literal_safe, _token


# ── Normalisering ─────────────────────────────────────────────────────────────

def _norm(s: str) -> str:
    """NFC + normaliser whitespace for robust matching mot originaltekst."""
    s = unicodedata.normalize("NFC", s)
    return re.sub(r"\s+", " ", s).strip()


def _apply_replacements(text: str, replacements: dict[str, str]) -> str:
    """To-fase erstatning: tokens forhindrer kollisjon mellom overlappende funn."""
    if not replacements:
        return text

    # Fase 1: erstatt med XML-kompatible private-use tokens (lengste-først).
    tokens: list[tuple[str, str]] = []
    sorted_items = sorted(replacements.items(), key=lambda x: len(x[0]), reverse=True)
    for i, (old, new) in enumerate(sorted_items):
        old = _clean_replacement_text(old)
        if not old:
            continue
        token = _token(i)
        text = _replace_literal_safe(text, old, token)
        norm_old = _norm(old)
        if norm_old != old:
            text = _replace_literal_safe(text, norm_old, token)
        tokens.append((token, _clean_replacement_text(new)))

    # Fase 2: tokens → endelige plassholdere
    for token, new in tokens:
        text = text.replace(token, new)
    return text


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _replace_in_para(para, replacements: dict[str, str]) -> None:
    """Erstatt tekst i ett avsnitt.  Merger runs slik at split-tekst matches."""
    if not para.runs:
        return
    full = "".join(r.text for r in para.runs)
    modified = _apply_replacements(full, replacements)
    if modified == full:
        return
    # Skriv tilbake: legg alt i første run, nullstill resten.
    # Første runs formatering (font, størrelse) beholdes for hele avsnittet.
    para.runs[0].text = modified
    for run in para.runs[1:]:
        run.text = ""


def _process_docx_paras(paragraphs, replacements: dict[str, str]) -> None:
    for para in paragraphs:
        _replace_in_para(para, replacements)


# ── Office/PDF annotations ───────────────────────────────────────────────────

_REL_NS = "{http://schemas.openxmlformats.org/package/2006/relationships}"
_CT_NS = "{http://schemas.openxmlformats.org/package/2006/content-types}"


def _rewrite_zip(path: Path, remove_file, transform_file) -> None:
    tmp = path.with_name(f"{path.name}.tmp")
    try:
        with zipfile.ZipFile(path, "r") as zin, zipfile.ZipFile(tmp, "w") as zout:
            for item in zin.infolist():
                name = item.filename
                if remove_file(name):
                    continue
                data = zin.read(item)
                data = transform_file(name, data)
                zout.writestr(item, data)
        os.replace(tmp, path)
    finally:
        tmp.unlink(missing_ok=True)


def _remove_elements_by_localname(data: bytes, localnames: set[str]) -> bytes:
    try:
        root = ET.fromstring(data)
    except ET.ParseError:
        return data
    parent_map = {child: parent for parent in root.iter() for child in parent}
    for elem in list(root.iter()):
        local = elem.tag.rsplit("}", 1)[-1] if "}" in elem.tag else elem.tag
        if local in localnames and elem in parent_map:
            parent_map[elem].remove(elem)
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _strip_relationships(data: bytes, markers: tuple[str, ...]) -> bytes:
    try:
        root = ET.fromstring(data)
    except ET.ParseError:
        return data
    for rel in list(root):
        type_attr = rel.attrib.get("Type", "").casefold()
        target_attr = rel.attrib.get("Target", "").casefold()
        if any(marker in type_attr or marker in target_attr for marker in markers):
            root.remove(rel)
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _strip_content_type_overrides(data: bytes, part_markers: tuple[str, ...]) -> bytes:
    try:
        root = ET.fromstring(data)
    except ET.ParseError:
        return data
    for elem in list(root):
        part_name = elem.attrib.get("PartName", "").casefold()
        if elem.tag == f"{_CT_NS}Override" and any(marker in part_name for marker in part_markers):
            root.remove(elem)
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def strip_docx_annotations(path: Path) -> None:
    """Fjern DOCX-kommentarer og synlige kommentarankere fra dokumentpakken."""
    comment_parts = ("word/comments", "word/people.xml")

    def remove_file(name: str) -> bool:
        lowered = name.casefold()
        return lowered.startswith(comment_parts)

    def transform_file(name: str, data: bytes) -> bytes:
        lowered = name.casefold()
        if lowered == "[content_types].xml":
            return _strip_content_type_overrides(data, ("/word/comments", "/word/people.xml"))
        if lowered.endswith(".rels"):
            return _strip_relationships(data, ("comments", "people"))
        if lowered.startswith("word/") and lowered.endswith(".xml"):
            return _remove_elements_by_localname(
                data,
                {"commentRangeStart", "commentRangeEnd", "commentReference"},
            )
        return data

    _rewrite_zip(path, remove_file, transform_file)


def strip_pptx_annotations(path: Path) -> None:
    """Fjern PPTX speaker notes og kommentar-deler fra dokumentpakken."""
    removed_parts = (
        "ppt/notesslides/",
        "ppt/notesmasters/",
        "ppt/comments/",
        "ppt/commentauthors.xml",
        "ppt/threadedcomments/",
    )

    def remove_file(name: str) -> bool:
        return name.casefold().startswith(removed_parts)

    def transform_file(name: str, data: bytes) -> bytes:
        lowered = name.casefold()
        markers = ("notesslide", "notesmaster", "comments", "commentauthors", "threadedcomments")
        if lowered == "[content_types].xml":
            return _strip_content_type_overrides(
                data,
                (
                    "/ppt/notesslides/",
                    "/ppt/notesmasters/",
                    "/ppt/comments/",
                    "/ppt/commentauthors.xml",
                    "/ppt/threadedcomments/",
                ),
            )
        if lowered.endswith(".rels"):
            return _strip_relationships(data, markers)
        return data

    _rewrite_zip(path, remove_file, transform_file)


def strip_xlsx_annotations(path: Path) -> None:
    """Fjern XLSX-kommentarer/threaded comments og VML-comment drawings."""
    removed_parts = (
        "xl/comments",
        "xl/threadedcomments/",
        "xl/persons/",
        "xl/drawings/vmldrawing",
    )

    def remove_file(name: str) -> bool:
        return name.casefold().startswith(removed_parts)

    def transform_file(name: str, data: bytes) -> bytes:
        lowered = name.casefold()
        markers = ("comments", "threadedcomments", "persons", "vmldrawing")
        if lowered == "[content_types].xml":
            return _strip_content_type_overrides(
                data,
                (
                    "/xl/comments",
                    "/xl/threadedcomments/",
                    "/xl/persons/",
                    "/xl/drawings/vmldrawing",
                ),
            )
        if lowered.endswith(".rels"):
            return _strip_relationships(data, markers)
        if lowered.startswith("xl/worksheets/") and lowered.endswith(".xml"):
            return _remove_elements_by_localname(data, {"legacyDrawing", "legacyDrawingHF"})
        return data

    _rewrite_zip(path, remove_file, transform_file)


def _has_existing_header_footer_definition(header_footer) -> bool:
    """Return True when python-docx can read this header/footer without creating one."""
    current = header_footer
    while current is not None:
        if current._has_definition:  # noqa: SLF001 - avoids creating default header/footer parts
            return True
        current = current._prior_headerfooter  # noqa: SLF001
    return False


def patch_docx(source: Path, replacements: dict[str, str], output: Path, *, strip_annotations: bool = False) -> None:
    from docx import Document  # type: ignore[import-untyped]

    doc = Document(str(source))

    # Brødtekst
    _process_docx_paras(doc.paragraphs, replacements)

    # Tabeller i brødtekst
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                _process_docx_paras(cell.paragraphs, replacements)

    # Topptekst, bunntekst og tabeller der
    for section in doc.sections:
        for hf in (section.header, section.footer):
            if not _has_existing_header_footer_definition(hf):
                continue
            _process_docx_paras(hf.paragraphs, replacements)
            for table in hf.tables:
                for row in table.rows:
                    for cell in row.cells:
                        _process_docx_paras(cell.paragraphs, replacements)

    doc.save(str(output))
    if strip_annotations:
        strip_docx_annotations(output)


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _iter_pptx_shapes(shapes):
    """Yield alle shapes rekursivt (inkl. grupperte shapes)."""
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):   # GroupShapes
            yield from _iter_pptx_shapes(shape.shapes)


def patch_pptx(source: Path, replacements: dict[str, str], output: Path, *, strip_annotations: bool = False) -> None:
    from pptx import Presentation  # type: ignore[import-untyped]

    prs = Presentation(str(source))

    for slide in prs.slides:
        for shape in _iter_pptx_shapes(slide.shapes):
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    _replace_in_para(para, replacements)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            _replace_in_para(para, replacements)

    prs.save(str(output))
    if strip_annotations:
        strip_pptx_annotations(output)


# ── XLSX ──────────────────────────────────────────────────────────────────────

def patch_xlsx(source: Path, replacements: dict[str, str], output: Path, *, strip_annotations: bool = False) -> None:
    from openpyxl import load_workbook  # type: ignore[import-untyped]

    wb = load_workbook(str(source))
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    cell.value = _apply_replacements(cell.value, replacements)
                if strip_annotations and getattr(cell, "comment", None) is not None:
                    cell.comment = None
    wb.save(str(output))
    if strip_annotations:
        strip_xlsx_annotations(output)


# ── PDF ───────────────────────────────────────────────────────────────────────

def redact_pdf(source: Path, replacements: dict[str, str], output: Path, *, strip_annotations: bool = False) -> None:
    """PDF-redaksjon med pymupdf.

    Erstatter sensitiv tekst med plassholder-tekst på hvit bakgrunn.
    Tekst som ikke finnes på siden hoppes over stille.
    """
    import fitz  # type: ignore[import-untyped]  # pymupdf

    doc = fitz.open(str(source))
    for page in doc:
        if strip_annotations:
            for annot in list(page.annots() or []):
                page.delete_annot(annot)
        for old, new in replacements.items():
            hits = page.search_for(old)
            for rect in hits:
                page.add_redact_annot(
                    rect,
                    text=new,
                    fontname="helv",
                    fontsize=9,
                    fill=(1, 1, 1),        # hvit bakgrunn
                    text_color=(0, 0, 0),  # svart tekst
                )
            # Prøv normalisert variant
            norm_old = _norm(old)
            if norm_old != old:
                for rect in page.search_for(norm_old):
                    page.add_redact_annot(
                        rect,
                        text=new,
                        fontname="helv",
                        fontsize=9,
                        fill=(1, 1, 1),
                        text_color=(0, 0, 0),
                    )
        # PyMuPDF exposes apply_redactions(); older examples sometimes use
        # apply_redacts(), which is not available in current builds.
        page.apply_redactions()

    doc.save(str(output), garbage=4, deflate=True)
    doc.close()


# ── Dispatcher ────────────────────────────────────────────────────────────────

SUPPORTED_PATCH_SUFFIXES = {".docx", ".pptx", ".xlsx", ".pdf"}


def patch_file(
    source: Path,
    replacements: dict[str, str],
    output: Path,
    *,
    strip_annotations: bool = False,
) -> None:
    """Kall riktig patcher basert på filsuffikset."""
    suffix = source.suffix.lower()
    if suffix == ".docx":
        patch_docx(source, replacements, output, strip_annotations=strip_annotations)
    elif suffix == ".pptx":
        patch_pptx(source, replacements, output, strip_annotations=strip_annotations)
    elif suffix == ".xlsx":
        patch_xlsx(source, replacements, output, strip_annotations=strip_annotations)
    elif suffix == ".pdf":
        redact_pdf(source, replacements, output, strip_annotations=strip_annotations)
    else:
        raise ValueError(f"Filformat ikke støttet for in-place anonymisering: {suffix}")
