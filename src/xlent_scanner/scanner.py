"""Tekstekstraksjon og scan-orkestrering."""
from __future__ import annotations

import html
import re
import tempfile
import time
import warnings
from collections.abc import Iterable
from contextvars import ContextVar
from pathlib import Path

import fitz  # type: ignore[import-untyped]  # fallback PDF-parser

_pdf_converter = None      # Docling DocumentConverter – lazy-initialisert ved første PDF-scan
_pdf_converter_ocr = None  # OCR-variant (do_ocr=True) – kun initialisert på eksplisitt forespørsel
_image_ocr_engine = None   # RapidOCR – lazy-initialisert ved første bilde-OCR
_DOCLING_TABLE_IMAGE_DEPRECATION = (
    r"This field is deprecated\. Use `generate_page_images=True` and call "
    r"`TableItem\.get_image\(\)` to extract table images from page images\."
)
_DOCLING_IMAGE_PLACEHOLDER_RE = re.compile(r"(?is)<!--\s*image\s*-->")
_extraction_metadata: ContextVar[dict[str, str]] = ContextVar("xlent_scanner_extraction_metadata", default={})

from xlent_scanner.detectors.clients import detect_clients
from xlent_scanner.detectors.creditcards import detect_creditcards
from xlent_scanner.detectors.financials import detect_financials
from xlent_scanner.detectors.iban import detect_iban
from xlent_scanner.detectors.keywords import detect_keywords
from xlent_scanner.detectors.ner_names import detect_names, get_load_error
from xlent_scanner.detectors.regex_da import detect_da_specific
from xlent_scanner.detectors.regex_de import detect_de_specific
from xlent_scanner.detectors.regex_en import detect_en_specific
from xlent_scanner.detectors.regex_es import detect_es_specific
from xlent_scanner.detectors.regex_extra import detect_extra
from xlent_scanner.detectors.regex_fr import detect_fr_specific
from xlent_scanner.detectors.regex_no import detect_no_specific, find_emails
from xlent_scanner.detectors.regex_sv import detect_sv_specific
from xlent_scanner.detectors.regex_url import detect_urls
from xlent_scanner.detectors.secrets import detect_secrets
from xlent_scanner.ignore import filter_findings, load_ignore_list
from xlent_scanner.language import resolve_language
from xlent_scanner.models import Finding, ScanResult  # noqa: F401
from xlent_scanner.risk import assess
from xlent_scanner.scan_categories import (
    category_enabled as _category_enabled,
    finding_matches_scan_categories as _finding_matches_scan_categories,
    normalise_scan_categories,
)
from xlent_scanner.suppression import capture_suppressed_findings
from xlent_scanner.whitelist import mark_whitelist_findings
from xlent_scanner.blacklist import detect_blacklist
from xlent_scanner.detectors.custom_patterns import detect_custom_patterns

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}
TEXT_WARNING_MIN_CHARS = 50

SUPPORTED_SUFFIXES = {
    ".pdf", ".docx", ".pptx", ".xlsx",
    ".md", ".txt", ".html",
    ".csv", ".eml", ".rtf", ".odt",
} | IMAGE_SUFFIXES

DEFAULT_FOLDER_MAX_FILES = 500
DEFAULT_FOLDER_MAX_DEPTH = 5
DEFAULT_EXCLUDED_DIRS = {
    ".git", ".hg", ".svn",
    ".mypy_cache", ".pytest_cache", ".ruff_cache",
    ".venv", "venv", "__pycache__",
    "node_modules",
    "build", "dist",
    "Library", "AppData",
}

_ignore_list: dict | None = None


_TECHNICAL_PROFILE_PATTERNS: tuple[tuple[re.Pattern[str], int], ...] = (
    (re.compile(r"(?i)\bdoi\b|doi\.org/|10\.\d{4,9}/"), 3),
    (re.compile(r"(?i)\b(?:isbn|issn)\b"), 2),
    (re.compile(r"(?i)\bet\s+al\."), 2),
    (re.compile(r"(?i)\b(?:references|bibliography)\b"), 1),
    (re.compile(r"(?i)\b(?:abstract|methodology|method|results|discussion)\b"), 1),
    (re.compile(r"(?i)\b(?:fig\.?|figure|table)\s*\d+\b"), 1),
    (re.compile(r"(?i)\b(?:frequency|amplitude|spectrum|axis|time\s+domain|phase|solver)\b"), 1),
    (re.compile(r"(?i)\b\d+(?:[.,]\d+)?\s*(?:hz|khz|mhz|ghz|ms|ns|mm|cm|km|m/s|db|°c)\b"), 1),
    (re.compile(r"\[[0-9]{1,3}\]"), 1),
)


def normalise_scan_profile(scan_profile: str | None) -> str:
    profile = str(scan_profile or "").strip().lower()
    return "technical" if profile in {"technical", "academic"} else "normal"


def looks_like_technical_or_academic_text(text: str) -> bool:
    """Heuristisk profilvalg for dokumenter med teknisk/akademisk støy.

    Terskelen er bevisst konservativ: ett svakt signal skal ikke endre profil,
    mens DOI/DOI-lenker alene er et sterkt akademisk signal.
    """
    sample = str(text or "")[:50_000]
    if not sample.strip():
        return False
    score = 0
    for pattern, weight in _TECHNICAL_PROFILE_PATTERNS:
        matches = pattern.findall(sample)
        if not matches:
            continue
        score += min(len(matches), 3) * weight
        if weight >= 3:
            return True
    return score >= 3


def resolve_scan_profile(scan_profile: str | None, text: str) -> str:
    profile = str(scan_profile or "normal").strip().lower()
    if profile == "auto":
        return "technical" if looks_like_technical_or_academic_text(text) else "normal"
    return normalise_scan_profile(profile)


def normalise_pdf_mode(pdf_mode: str | None) -> str:
    mode = str(pdf_mode or "fast").strip().lower()
    return mode if mode in {"fast", "auto", "advanced"} else "fast"


def _set_extraction_metadata(**values: str) -> None:
    _extraction_metadata.set({key: value for key, value in values.items() if value})


def _current_extraction_metadata() -> dict[str, str]:
    return dict(_extraction_metadata.get() or {})


def reset_ignore_cache() -> None:
    global _ignore_list
    _ignore_list = None


def _get_ignore_list() -> dict:
    global _ignore_list
    if _ignore_list is None:
        _ignore_list = load_ignore_list()
    return _ignore_list


def _read_text_file(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="ignore")


def _ignore_docling_table_image_deprecation() -> None:
    warnings.filterwarnings(
        "ignore",
        message=_DOCLING_TABLE_IMAGE_DEPRECATION,
        category=DeprecationWarning,
        module=r"docling\..*",
    )


def _build_pdf_converter(ocr: bool):
    from docling.document_converter import DocumentConverter, PdfFormatOption
    from docling.datamodel.pipeline_options import PdfPipelineOptions
    from docling.datamodel.base_models import InputFormat

    opts = PdfPipelineOptions()
    opts.do_ocr = ocr
    # Docling laster layout-modell automatisk ved første konvertering.
    # Docling 2.x leser et deprecated internt bildefelt selv om vi ikke bruker det.
    # Filtrer kun denne kjente tredjeparts-warningen uten å endre PDF-oppførsel.
    with warnings.catch_warnings():
        _ignore_docling_table_image_deprecation()
        return DocumentConverter(
            format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=opts)}
        )


def _get_pdf_converter(ocr: bool = False):
    """Returnerer (og cacher) en Docling DocumentConverter for PDF-parsing.

    OCR-varianten caches separat — den drar inn OCR-motoren (EasyOCR) og
    skal kun initialiseres når brukeren eksplisitt ber om OCR.
    Raises RuntimeError hvis Docling eller en av dens avhengigheter ikke er tilgjengelig.
    """
    global _pdf_converter, _pdf_converter_ocr
    if ocr:
        if _pdf_converter_ocr is None:
            _pdf_converter_ocr = _build_pdf_converter(ocr=True)
        return _pdf_converter_ocr
    if _pdf_converter is None:
        _pdf_converter = _build_pdf_converter(ocr=False)
    return _pdf_converter


def _extract_text_pdf_fitz(path: Path) -> str:
    """Fallback: enkel tekstekstraksjon via PyMuPDF."""
    chunks: list[str] = []
    with fitz.open(path) as doc:
        for page in doc:
            chunks.append(page.get_text("text"))
    return "\n".join(chunks).strip()


def _pdf_looks_image_based(path: Path) -> bool:
    """Returnerer True når PDF-en har bilder, men lite/ingen ekte PDF-tekst.

    Dette brukes kun som en OCR-hint. Docling kan returnere Markdown-markører som
    ``<!-- image -->`` for bildebaserte PDF-er; slike markører er ikke reell tekst
    og skal ikke hindre at GUI tilbyr OCR.
    """
    try:
        with fitz.open(path) as doc:
            if doc.page_count == 0:
                return False
            image_count = 0
            embedded_text_len = 0
            for page in doc:
                image_count += len(page.get_images(full=True))
                embedded_text_len += len(page.get_text("text").strip())
            return image_count > 0 and embedded_text_len < TEXT_WARNING_MIN_CHARS
    except Exception:
        return False


def _meaningful_extracted_text(text: str) -> str:
    """Fjerner parser-placeholders før vi vurderer om et dokument har tekst."""
    return _DOCLING_IMAGE_PLACEHOLDER_RE.sub(" ", text).strip()


_TABLE_KEYWORD_RE = re.compile(
    r"(?i)\b("
    r"unit\s+price|quantity|qty|description|invoice|subtotal|total|vat|amount|"
    r"table\s+\d+|kolonne|beløp|sum|mva|antall|enhetspris"
    r")\b"
)
_MULTISPACE_COLUMN_RE = re.compile(r"\S\s{2,}\S")
_NUMERIC_CELL_RE = re.compile(r"(?:\b\d{1,3}(?:[ .,\t]\d{3})+(?:[,.]\d{1,2})?\b|\b\d+[,.]\d{2}\b|\b\d+\s*(?:kr|nok|eur|usd|%|\$|€)\b)", re.I)


def _looks_like_table_text(text: str) -> bool:
    """Heuristikk for PDF-tekst som trolig representerer tabeller/kolonner.

    Brukes kun i PDF auto-modus som signal for å prøve Docling. Terskelen er
    konservativ for å unngå at vanlige avsnitt sender alle PDF-er gjennom
    tyngre layout-parser.
    """
    lines = [line.strip() for line in str(text or "").splitlines() if line.strip()]
    if len(lines) < 3:
        return False

    column_like = 0
    numeric_dense = 0
    keyword_hits = 0
    for line in lines[:300]:
        if _MULTISPACE_COLUMN_RE.search(line) or line.count("\t") >= 2 or line.count("|") >= 2:
            column_like += 1
        if len(_NUMERIC_CELL_RE.findall(line)) >= 2:
            numeric_dense += 1
        if _TABLE_KEYWORD_RE.search(line):
            keyword_hits += 1

    return (
        column_like >= 4
        or (column_like >= 2 and (numeric_dense >= 2 or keyword_hits >= 2))
        or (numeric_dense >= 4 and keyword_hits >= 1)
    )


def _extract_text_pdf(path: Path, ocr: bool = False, pdf_mode: str = "fast") -> str:
    """Ekstraherer tekst fra PDF.

    Standardmodus er rask PyMuPDF. Docling brukes bare i advanced/auto eller OCR.

    Med ocr=True kjøres Docling med OCR-motor aktivert (for bilde-PDF-er).
    Da er PyMuPDF-fallback meningsløs (den gir samme tomme resultat), så
    feil propageres som RuntimeError med forklarende melding i stedet.
    """
    if ocr:
        _set_extraction_metadata(scan_strategy="advanced", scan_strategy_reason="ocr")
        try:
            with warnings.catch_warnings():
                _ignore_docling_table_image_deprecation()
                converter = _get_pdf_converter(ocr=True)
                return converter.convert(str(path)).document.export_to_markdown()
        except Exception as exc:
            raise RuntimeError(
                "OCR er ikke tilgjengelig i denne installasjonen "
                f"({type(exc).__name__}: {exc}). "
                "OCR krever Docling med OCR-motor — kjør fra kildekode "
                "med «uv sync» hvis pakken mangler den."
            ) from exc
    mode = normalise_pdf_mode(pdf_mode)
    fitz_text = _extract_text_pdf_fitz(path)
    if mode == "fast":
        _set_extraction_metadata(scan_strategy="fast", scan_strategy_reason="explicit_fast")
        return fitz_text
    meaningful_fitz_len = len(_meaningful_extracted_text(fitz_text))
    table_like = _looks_like_table_text(fitz_text)
    if mode == "auto":
        if meaningful_fitz_len >= TEXT_WARNING_MIN_CHARS and not table_like:
            _set_extraction_metadata(scan_strategy="fast", scan_strategy_reason="auto_fast")
            return fitz_text
        docling_reason = "table_layout" if table_like else "little_text"
    else:
        docling_reason = "explicit_advanced"

    try:
        with warnings.catch_warnings():
            _ignore_docling_table_image_deprecation()
            converter = _get_pdf_converter()
            docling_text = converter.convert(str(path)).document.export_to_markdown()
    except Exception:
        _set_extraction_metadata(scan_strategy="fast", scan_strategy_reason=f"{docling_reason}_fallback")
        return fitz_text

    if len(_meaningful_extracted_text(docling_text)) > len(_meaningful_extracted_text(fitz_text)):
        _set_extraction_metadata(scan_strategy="advanced", scan_strategy_reason=docling_reason)
        return docling_text
    _set_extraction_metadata(scan_strategy="fast", scan_strategy_reason=f"{docling_reason}_not_better")
    return fitz_text


def _get_image_ocr_engine():
    """Returnerer en cached RapidOCR-motor for rene bildefiler."""
    global _image_ocr_engine
    if _image_ocr_engine is None:
        try:
            import onnxruntime  # noqa: F401, PLC0415
            from rapidocr import RapidOCR  # noqa: PLC0415
        except Exception as exc:
            raise RuntimeError(
                "OCR-motor mangler. Bildefiler krever rapidocr + onnxruntime."
            ) from exc
        _image_ocr_engine = RapidOCR()
    return _image_ocr_engine


def _image_to_temp_pdf(path: Path) -> Path:
    """Konverter en bildefil til midlertidig én-sides PDF for Docling OCR."""
    with fitz.open(path) as image_doc:
        pdf_bytes = image_doc.convert_to_pdf()
    with tempfile.NamedTemporaryFile(prefix="xlent-image-ocr-", suffix=".pdf", delete=False) as handle:
        handle.write(pdf_bytes)
        return Path(handle.name)


def _extract_text_image(path: Path, ocr: bool = False, pdf_mode: str = "fast") -> str:
    """Ekstraherer tekst fra PNG/JPG/TIFF/WebP via OCR.

    Uten eksplisitt OCR returnerer vi tom tekst slik at GUI kan vise samme
    OCR-tilbud som for bildebaserte PDF-er.

    Avansert scan-modus konverterer bildet til midlertidig PDF og kjører
    Docling OCR. Det gir bedre sjanse for å bevare layout/struktur enn ren
    RapidOCR-linjetekst.
    """
    if normalise_pdf_mode(pdf_mode) == "advanced":
        tmp_pdf: Path | None = None
        try:
            tmp_pdf = _image_to_temp_pdf(path)
            return _extract_text_pdf(tmp_pdf, ocr=True, pdf_mode="advanced")
        except Exception as exc:
            raise RuntimeError(
                "Avansert OCR av bildefil feilet. Kontroller at Docling OCR "
                f"er tilgjengelig ({type(exc).__name__}: {exc})."
            ) from exc
        finally:
            if tmp_pdf is not None:
                try:
                    tmp_pdf.unlink(missing_ok=True)
                except OSError:
                    pass
    if not ocr:
        _set_extraction_metadata(scan_strategy="fast", scan_strategy_reason="image_no_ocr")
        return ""
    _set_extraction_metadata(scan_strategy="fast", scan_strategy_reason="ocr")
    try:
        result = _get_image_ocr_engine()(path)
    except Exception as exc:
        raise RuntimeError(
            "OCR av bildefil feilet. Kontroller at full/OCR-build eller "
            f"kildeinstallasjon har OCR-avhengigheter ({type(exc).__name__}: {exc})."
        ) from exc
    texts = getattr(result, "txts", None) or ()
    return "\n".join(str(text).strip() for text in texts if str(text).strip())


def _table_to_markdown(table) -> str:
    """Konverterer en python-docx- eller python-pptx-tabell til Markdown-format.

    Bevarer kolonnerelasjonene slik at AI og regelbasert skanner forstår
    at «30» er en verdi i «Cost (NOK)»-kolonnen, ikke et løst tall.
    Duplikate celler fra sammenslåtte celler (merged) dedupliseres per rad.
    """
    rows_text: list[list[str]] = []
    for row in table.rows:
        # python-docx gjentar samme _tc-objekt for sammenslåtte celler.
        # Bruk XML-element-id for å deduplisere, ikke tekstinnhold –
        # ellers mister vi reelle celler som tilfeldigvis har samme tekst.
        seen_tc: set[int] = set()
        cells: list[str] = []
        for cell in row.cells:
            tc_id = id(getattr(cell, "_tc", cell))
            if tc_id not in seen_tc:
                seen_tc.add(tc_id)
                cells.append(cell.text.strip())
        if any(c for c in cells):
            rows_text.append(cells)

    if not rows_text:
        return ""

    # Normaliser kolonnebredde
    max_cols = max(len(r) for r in rows_text)
    lines: list[str] = []
    for i, row in enumerate(rows_text):
        # Padder kortere rader
        padded = row + [""] * (max_cols - len(row))
        lines.append(" | ".join(padded))
        if i == 0 and len(rows_text) > 1:
            lines.append(" | ".join(["---"] * max_cols))
    return "\n".join(lines)


def _extract_text_docx(path: Path) -> str:
    from docx import Document as DocxDocument  # lazy import for packaged stability

    doc = DocxDocument(str(path))
    parts: list[str] = []

    # Kombiner avsnitt og tabeller i dokumentets rekkefølge ved å iterere
    # over XML-barna direkte, slik at tabeller dukker opp på riktig sted
    # i teksten (f.eks. etter «Here is my budget:»-setningen).
    from docx.oxml.ns import qn  # noqa: PLC0415
    body = doc.element.body
    tbl_idx = 0
    para_idx = 0
    for child in body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "p":
            if para_idx < len(doc.paragraphs):
                txt = doc.paragraphs[para_idx].text
                if txt:
                    parts.append(txt)
                para_idx += 1
        elif tag == "tbl":
            if tbl_idx < len(doc.tables):
                md = _table_to_markdown(doc.tables[tbl_idx])
                if md:
                    parts.append(md)
                tbl_idx += 1

    # Fall-back: legg til eventuelle urelaterte tabeller som ikke ble funnet i body
    for i in range(tbl_idx, len(doc.tables)):
        md = _table_to_markdown(doc.tables[i])
        if md:
            parts.append(md)

    return "\n".join(parts).strip()


def _extract_text_pptx(path: Path) -> str:
    from pptx import Presentation  # lazy import for packaged stability

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                md = _table_to_markdown(shape.table)
                if md:
                    parts.append(md)
        try:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                parts.append(notes)
        except Exception:
            pass
    return "\n".join(parts).strip()


def _extract_text_xlsx(path: Path) -> str:
    from openpyxl import load_workbook  # lazy import for packaged stability

    wb = load_workbook(filename=str(path), data_only=True, read_only=True)
    try:
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"[{ws.title}]")
            for row in ws.iter_rows(values_only=True):
                vals = [str(v) for v in row if v is not None and str(v).strip()]
                if vals:
                    parts.append(" | ".join(vals))
        return "\n".join(parts).strip()
    finally:
        wb.close()


def _extract_text_csv(path: Path) -> str:
    """Leser CSV og returnerer celleinnhold rad for rad."""
    import csv as _csv
    content = _read_text_file(path)
    try:
        reader = _csv.reader(content.splitlines())
        parts = [" | ".join(cell for cell in row if cell.strip()) for row in reader]
        return "\n".join(p for p in parts if p).strip()
    except Exception:
        return content.strip()


def _extract_text_eml(path: Path) -> str:
    """Ekstraherer tekst fra e-post (.eml) via stdlib email-modulen."""
    import email as _email
    from email import policy as _policy
    content = _read_text_file(path)
    msg = _email.message_from_string(content, policy=_policy.default)
    parts: list[str] = []
    # Relevante headers
    for header in ("From", "To", "Cc", "Bcc", "Subject", "Date", "Reply-To"):
        val = msg.get(header, "")
        if val:
            parts.append(f"{header}: {val}")
    # Brødtekst
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                try:
                    parts.append(str(part.get_content()))
                except Exception:
                    pass
    else:
        try:
            parts.append(str(msg.get_content()))
        except Exception:
            pass
    return "\n".join(parts).strip()


# Regex for å strippe RTF-kontrollkoder (enkel men presis for vanlig tekst-RTF)
_RTF_CTRL = re.compile(
    r"\\(?:[a-z]+(?:-?\d+)? ?|\'[0-9a-fA-F]{2}|[^a-z\s])"
    r"|[{}]",
    re.IGNORECASE,
)


def _extract_text_rtf(path: Path) -> str:
    """Enkel RTF-ekstraksjon: fjerner kontrollkoder, beholder løpende tekst."""
    content = _read_text_file(path)
    if not content.lstrip().startswith("{\\rtf"):
        return content.strip()
    text = _RTF_CTRL.sub("", content)
    return re.sub(r"\s+", " ", text).strip()


def _extract_text_odt(path: Path) -> str:
    """Ekstraherer tekst fra ODT/ODS/ODP via zipfile + XML-parsing."""
    import xml.etree.ElementTree as _ET
    import zipfile as _zf

    with _zf.ZipFile(str(path), "r") as zf:
        if "content.xml" not in zf.namelist():
            return ""
        with zf.open("content.xml") as f:
            root = _ET.parse(f).getroot()

    parts: list[str] = []
    for elem in root.iter():
        if elem.text and elem.text.strip():
            parts.append(elem.text.strip())
        if elem.tail and elem.tail.strip():
            parts.append(elem.tail.strip())
    return "\n".join(parts).strip()


def _extract_text_html(path: Path) -> str:
    content = _read_text_file(path)
    content = re.sub(r"(?is)<script.*?>.*?</script>", " ", content)
    content = re.sub(r"(?is)<style.*?>.*?</style>", " ", content)
    content = re.sub(r"(?s)<[^>]+>", " ", content)
    content = html.unescape(content)
    content = re.sub(r"\s+", " ", content)
    return content.strip()


def extract_text(path: Path, ocr: bool = False, pdf_mode: str = "fast") -> str:
    suffix = path.suffix.lower()
    if suffix in IMAGE_SUFFIXES:
        return _extract_text_image(path, ocr=ocr, pdf_mode=pdf_mode)
    if suffix == ".pdf":
        return _extract_text_pdf(path, ocr=ocr, pdf_mode=pdf_mode)
    if suffix == ".docx":
        return _extract_text_docx(path)
    if suffix == ".pptx":
        return _extract_text_pptx(path)
    if suffix == ".xlsx":
        return _extract_text_xlsx(path)
    if suffix in {".md", ".txt"}:
        return _read_text_file(path).strip()
    if suffix == ".html":
        return _extract_text_html(path)
    if suffix == ".csv":
        return _extract_text_csv(path)
    if suffix == ".eml":
        return _extract_text_eml(path)
    if suffix == ".rtf":
        return _extract_text_rtf(path)
    if suffix == ".odt":
        return _extract_text_odt(path)
    raise ValueError(f"Filtype ikke støttet: {suffix}")


def _run_detectors(
    text: str,
    lang: str,
    ignore_xlent: bool = False,
    scan_profile: str = "normal",
    categories: Iterable[str] | None = None,
) -> tuple[list[Finding], bool, dict]:
    """Kjør alle detektorer og returner funn samt om skannen ble degradert."""
    findings: list[Finding] = []
    detector_errors: list[str] = []
    timings: dict[str, float] = {}

    scan_profile = resolve_scan_profile(scan_profile, text)
    selected_categories = normalise_scan_categories(categories)

    def _run(fn, *args, **kwargs):
        t0 = time.perf_counter()
        try:
            findings.extend(fn(*args, **kwargs))
        except Exception as exc:
            detector_errors.append(f"{fn.__name__}: {type(exc).__name__}: {exc}")
        finally:
            timings[fn.__name__] = round(timings.get(fn.__name__, 0.0) + time.perf_counter() - t0, 4)

    if _category_enabled(selected_categories, "hemmeligheter"):
        _run(detect_keywords, text)
    if _category_enabled(selected_categories, "hemmeligheter"):
        _run(detect_secrets, text)
    if _category_enabled(selected_categories, "epost"):
        _run(find_emails, text)
    if _category_enabled(selected_categories, "nettadresse"):
        _run(detect_urls, text)
    if lang in ("nb", "en", "da") and _category_enabled(
        selected_categories,
        "id",
        "klient",
        "konto",
        "telefon",
    ):
        _run(detect_no_specific, text, scan_profile=scan_profile)
    if lang in ("sv", "en") and _category_enabled(
        selected_categories,
        "id",
        "klient",
        "konto",
        "telefon",
    ):
        _run(detect_sv_specific, text)
    if lang == "da" and _category_enabled(selected_categories, "id"):
        _run(detect_da_specific, text)
    if lang in ("en", "de", "fr", "es") and _category_enabled(selected_categories, "id", "telefon"):
        _run(detect_en_specific, text)
    if lang == "de" and _category_enabled(selected_categories, "id", "telefon"):
        _run(detect_de_specific, text)
    if lang == "fr" and _category_enabled(selected_categories, "id", "telefon"):
        _run(detect_fr_specific, text)
    if lang == "es" and _category_enabled(selected_categories, "id", "telefon"):
        _run(detect_es_specific, text)
    if _category_enabled(selected_categories, "konto"):
        _run(detect_iban, text)
    if _category_enabled(selected_categories, "konto"):
        _run(detect_creditcards, text)
    if _category_enabled(selected_categories, "finansielt"):
        _run(detect_financials, text)
    if _category_enabled(selected_categories, "klient"):
        _run(detect_clients, text)
    if _category_enabled(
        selected_categories,
        "id",
        "konto",
        "nettadresse",
        "telefon",
        "adresse",
        "finansielt",
        "medisinsk",
        "hemmeligheter",
    ):
        _run(detect_extra, text)
    _run(detect_custom_patterns, text)
    if _category_enabled(selected_categories, "navn"):
        _run(detect_names, text, lang, scan_profile=scan_profile)

    if ignore_xlent:
        findings = filter_findings(findings, _get_ignore_list())

    findings = [f for f in findings if _finding_matches_scan_categories(f, selected_categories)]
    findings = mark_whitelist_findings(findings)
    _run(detect_blacklist, text)

    ner_err = get_load_error(lang) if _category_enabled(selected_categories, "navn") else None
    if ner_err:
        findings.append(Finding(category="⚠ NER ikke tilgjengelig", text=ner_err, context=""))
    for det_err in detector_errors:
        findings.append(Finding(category="⚠ Detektor-feil", text=det_err, context=""))

    return findings, bool(ner_err or detector_errors), timings


def scan_text(
    text: str,
    language: str = "auto",
    source_name: str = "Innlimt tekst",
    scan_profile: str = "normal",
    categories: Iterable[str] | None = None,
) -> ScanResult:
    """Skann ren tekst direkte (uten filekstraksjon). Brukes for utklippstavle/paste."""
    total_t0 = time.perf_counter()
    if not text.strip():
        return ScanResult(
            file_name=source_name, file_size=0, text_length=0, text_preview="",
            error="Ingen tekst å skanne.", scan_status="failed",
        )
    lang_t0 = time.perf_counter()
    lang = resolve_language(language, text)
    language_seconds = round(time.perf_counter() - lang_t0, 4)
    scan_profile = resolve_scan_profile(scan_profile, text)
    selected_categories = normalise_scan_categories(categories)
    with capture_suppressed_findings() as suppressed:
        detector_t0 = time.perf_counter()
        findings, degraded, detector_timings = _run_detectors(
            text,
            lang,
            scan_profile=scan_profile,
            categories=selected_categories,
        )
        detector_seconds = round(time.perf_counter() - detector_t0, 4)

    result = ScanResult(
        file_name=source_name,
        file_size=len(text.encode("utf-8")),
        text_length=len(text),
        text_preview=text,
        findings=findings,
        suppressed_findings=[
            f for f in suppressed if _finding_matches_scan_categories(f, selected_categories)
        ],
        original_text=text,
        language=lang,
        scan_status="partial" if degraded else "success",
        scan_timings={
            "language_seconds": language_seconds,
            "detectors_seconds": detector_seconds,
            "detectors": detector_timings,
            "total_seconds": round(time.perf_counter() - total_t0, 4),
            "scan_profile": scan_profile,
        },
    )
    return assess(result)


def _normalise_scan_limit(value: int | str | None, default: int, minimum: int, maximum: int) -> int:
    try:
        parsed = int(value) if value is not None else default
    except (TypeError, ValueError):
        parsed = default
    return max(minimum, min(parsed, maximum))


def _is_excluded_scan_dir(path: Path, excluded_dirs: set[str]) -> bool:
    name = path.name
    excluded_lower = {item.lower() for item in excluded_dirs}
    return path.is_symlink() or name.startswith(".") or name in excluded_dirs or name.lower() in excluded_lower


def build_folder_scan_plan(
    folder: str | Path,
    recursive: bool = False,
    max_files: int = DEFAULT_FOLDER_MAX_FILES,
    max_depth: int = DEFAULT_FOLDER_MAX_DEPTH,
    excluded_dirs: set[str] | None = None,
) -> dict:
    """Finn støttede filer i en mappe uten å lese filinnhold."""
    root = Path(folder)
    if not root.is_dir():
        raise ValueError(f"Ikke en mappe: {folder}")

    max_files = _normalise_scan_limit(max_files, DEFAULT_FOLDER_MAX_FILES, 1, 10_000)
    max_depth = _normalise_scan_limit(max_depth, DEFAULT_FOLDER_MAX_DEPTH, 0, 50)
    excluded = set(DEFAULT_EXCLUDED_DIRS if excluded_dirs is None else excluded_dirs)

    files: list[Path] = []
    folder_count = 0
    truncated = False

    def add_file(path: Path) -> bool:
        nonlocal truncated
        if path.suffix.lower() not in SUPPORTED_SUFFIXES:
            return True
        if len(files) >= max_files:
            truncated = True
            return False
        files.append(path)
        return True

    if not recursive:
        folder_count = 1
        for child in sorted(root.iterdir(), key=lambda p: p.name.lower()):
            if child.is_file() and not add_file(child):
                break
        files.sort(key=lambda p: p.name.lower())
    else:
        stack: list[tuple[Path, int]] = [(root, 0)]
        stop = False
        while stack and not stop:
            current, depth = stack.pop()
            folder_count += 1
            try:
                entries = sorted(current.iterdir(), key=lambda p: (not p.is_dir(), p.name.lower()))
            except OSError:
                continue
            child_dirs: list[tuple[Path, int]] = []
            for entry in entries:
                if entry.is_dir():
                    if depth < max_depth and not _is_excluded_scan_dir(entry, excluded):
                        child_dirs.append((entry, depth + 1))
                    continue
                if entry.is_file() and not add_file(entry):
                    stop = True
                    break
            stack.extend(reversed(child_dirs))
        files.sort(key=lambda p: str(p.relative_to(root)).lower())

    samples = [str(path.relative_to(root)) for path in files[:10]]
    return {
        "folder": str(root),
        "recursive": bool(recursive),
        "max_files": max_files,
        "max_depth": max_depth,
        "folder_count": folder_count,
        "file_count": len(files),
        "truncated": truncated,
        "files": files,
        "samples": samples,
        "excluded_dirs": sorted(excluded),
    }


def scan_folder(
    folder: str | Path,
    ignore_xlent: bool = False,
    language: str = "auto",
    max_files: int = DEFAULT_FOLDER_MAX_FILES,
    recursive: bool = False,
    max_depth: int = DEFAULT_FOLDER_MAX_DEPTH,
    scan_profile: str = "normal",
    categories: Iterable[str] | None = None,
    pdf_mode: str = "fast",
) -> list[ScanResult]:
    """Skann alle støttede filer i en mappe. Returnerer resultater sortert etter risikonivå."""
    plan = build_folder_scan_plan(
        folder,
        recursive=recursive,
        max_files=max_files,
        max_depth=max_depth,
    )
    root = Path(folder)
    results = []
    for f in plan["files"]:
        try:
            result = scan_file(
                f,
                ignore_xlent=ignore_xlent,
                language=language,
                scan_profile=scan_profile,
                categories=categories,
                pdf_mode=pdf_mode,
            )
        except TypeError as exc:
            if "unexpected keyword argument" not in str(exc):
                raise
            result = scan_file(f, ignore_xlent=ignore_xlent, language=language)
        result.relative_path = str(Path(f).relative_to(root))
        result.source_path = str(f)
        results.append(result)
    level_order = {"svart": 3, "rød": 2, "gul": 1, "grønn": 0}
    results.sort(key=lambda r: (-level_order.get(r.risk_level, 0), r.relative_path.lower()))
    return results


def scan_file(
    path: str | Path,
    ignore_xlent: bool = False,
    language: str = "auto",
    ocr: bool = False,
    scan_profile: str = "normal",
    categories: Iterable[str] | None = None,
    pdf_mode: str = "fast",
) -> ScanResult:
    total_t0 = time.perf_counter()
    p = Path(path)
    if not p.exists():
        return ScanResult(
            file_name=p.name, file_size=0, text_length=0, text_preview="",
            error=f"Fil ikke funnet: {p}",
            scan_status="failed",
        )
    suffix = p.suffix.lower()
    scan_mode = normalise_pdf_mode(pdf_mode)
    _set_extraction_metadata()
    advanced_image_ocr = suffix in IMAGE_SUFFIXES and scan_mode == "advanced"
    if suffix not in SUPPORTED_SUFFIXES:
        return ScanResult(
            file_name=p.name, file_size=p.stat().st_size,
            text_length=0, text_preview="",
            error=f"Filtype ikke støttet: {suffix}",
            scan_status="failed",
        )
    try:
        try:
            extract_t0 = time.perf_counter()
            text = extract_text(p, ocr=ocr, pdf_mode=scan_mode)
            extract_seconds = round(time.perf_counter() - extract_t0, 4)
        except TypeError as exc:
            # Flere tester og tredjepartsintegrasjoner monkeypatcher extract_text(path).
            # Behold bakoverkompatibilitet med eldre signaturer uten pdf_mode.
            if "unexpected keyword argument" not in str(exc):
                raise
            extract_t0 = time.perf_counter()
            try:
                text = extract_text(p, ocr=ocr)
            except TypeError as exc2:
                if ocr or "unexpected keyword argument" not in str(exc2):
                    raise
                text = extract_text(p)
            extract_seconds = round(time.perf_counter() - extract_t0, 4)
    except Exception as exc:
        msg = str(exc)
        if any(k in msg.lower() for k in ("illegal character", "not well-formed", "invalid token", "xml", "expat")):
            friendly = (
                "Filen inneholder ugyldige tegn og kunne ikke leses automatisk. "
                "Dette skjer gjerne med innscannede eller redigerte PPTX/DOCX-filer. "
                "Prøv å åpne filen og lagre den på nytt, eller konverter til PDF."
            )
        else:
            friendly = f"Klarte ikke å lese fil: {msg}"
        return ScanResult(
            file_name=p.name, file_size=p.stat().st_size,
            text_length=0, text_preview="",
            error=friendly,
            scan_status="failed",
        )

    extraction_metadata = _current_extraction_metadata()
    preview = text

    warning: str | None = None
    warning_code: str | None = None
    meaningful_text = _meaningful_extracted_text(text)
    pdf_image_based = suffix == ".pdf" and not ocr and _pdf_looks_image_based(p)
    if not meaningful_text:
        warning_code = "no_text_extracted"
        warning = (
            "Ingen tekst ble funnet i dokumentet. "
            "Filen kan være en bildebasert eller innskannet PDF/dokument. "
            "Innhold i bilder kan ikke sjekkes sikkert uten OCR. "
            "Bruk en tekstbasert versjon av filen, eller kjør OCR før scanning."
        )
    elif len(meaningful_text) < TEXT_WARNING_MIN_CHARS or (
        pdf_image_based
        and _DOCLING_IMAGE_PLACEHOLDER_RE.search(text)
        and len(meaningful_text) < TEXT_WARNING_MIN_CHARS * 4
    ):
        warning_code = "little_text_extracted"
        warning = (
            "Lite eller ingen tekst ble funnet i dokumentet. "
            "Filen kan bestå av innscannede bilder (bilde-PDF) "
            "og innholdet kan ikke sjekkes automatisk. "
            "Vurder å bruke en tekstbasert versjon av filen."
        )

    lang_t0 = time.perf_counter()
    lang = resolve_language(language, text)
    language_seconds = round(time.perf_counter() - lang_t0, 4)
    scan_profile = resolve_scan_profile(scan_profile, text)
    selected_categories = normalise_scan_categories(categories)

    with capture_suppressed_findings() as suppressed:
        detector_t0 = time.perf_counter()
        findings, degraded, detector_timings = _run_detectors(
            text,
            lang,
            ignore_xlent=ignore_xlent,
            scan_profile=scan_profile,
            categories=selected_categories,
        )
        detector_seconds = round(time.perf_counter() - detector_t0, 4)

    result = ScanResult(
        file_name=p.name,
        file_size=p.stat().st_size,
        text_length=len(text),
        text_preview=preview,
        findings=findings,
        suppressed_findings=[
            f for f in suppressed if _finding_matches_scan_categories(f, selected_categories)
        ],
        original_text=text,
        language=lang,
        warning=warning,
        warning_code=warning_code,
        ocr_used=bool(ocr or advanced_image_ocr),
        scan_status="partial" if warning_code or degraded else "success",
        scan_timings={
            "extract_seconds": extract_seconds,
            "language_seconds": language_seconds,
            "detectors_seconds": detector_seconds,
            "detectors": detector_timings,
            "total_seconds": round(time.perf_counter() - total_t0, 4),
            "scan_profile": scan_profile,
            "pdf_mode": scan_mode if suffix == ".pdf" else "",
            "scan_mode": scan_mode if suffix == ".pdf" or suffix in IMAGE_SUFFIXES else "",
            "scan_strategy": extraction_metadata.get("scan_strategy", ""),
            "scan_strategy_reason": extraction_metadata.get("scan_strategy_reason", ""),
        },
    )
    return assess(result)
